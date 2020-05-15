import nbformat
import jupytext
from jupytext.cli import jupytext as jupytext_cli
from pathlib import Path
from collections import defaultdict
from pptx import Presentation
import pprint
import string
import itertools
import click

pp = pprint.PrettyPrinter(indent=2)


def find_text(pptx_file):
    """
    given the path to a pptx file, return a dictionary
    with keys as integer slide numbers and values as a flat
    list of text strings on each slide
    """
    pptx_file = Path(pptx_file)
    the_pres = Presentation(pptx_file)
    slide_dict = {}
    for counter, slide in enumerate(the_pres.slides):
        text_list = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                before = shape.text.split("\n")
                after = [item.strip() for item in before if len(item.strip()) > 0]
                text_list.append(after)
        #
        # flatten this nested list using intertools.chain
        #
        slide_dict[counter] = {"text_list": list(itertools.chain(*text_list))}
    return slide_dict


def sort_fun(file_name):
    """
    sort image names of the form 
    week01_C1_Introduction2018-slide_9-006.png
    by image number (i.e. 006)
    """
    file_name = Path(file_name)
    file_name = file_name.stem
    name, slidenum, imagenum = str(file_name).split("-")
    return int(imagenum)


def find_images(media_dir, slide_dict):
    """
    get the png filenames in a media folder written by image_extract.py
    and return them as a dictionary indexed by integer slide number
    """
    media_dir = Path(media_dir)
    image_files = media_dir.glob("*.png")
    image_dict = defaultdict(list)
    for a_file in image_files:
        # week01_C1_Introduction2018-slide_9-006
        the_name, slide, image_num = a_file.stem.split("-")
        _, slide_num = slide.split("_")
        slide_num, image_num = int(slide_num), int(image_num)
        image_dict[slide_num].append(str(a_file))
    for slide_num, image_list in image_dict.items():
        image_list = sorted(image_list, key=sort_fun)
        print(f"here is image list: {image_list}")
        slide_dict[slide_num]["image_list"] = image_list
    return slide_dict


@click.command()
@click.argument("pptx_file", type=str, nargs=1)
@click.argument("media_directory", type=str, nargs=1)
def main(pptx_file, media_directory):
    """\b
    usage: python build_notebook.py pptx_file media_directory
    pptx_file: powerpoint file to convert
    media_directory: folder of png files produced by running image_extract.py on pptx_file
    """
    pptx_file = Path(pptx_file)
    text_dict = find_text(pptx_file)
    slide_dict = find_images(media_directory, text_dict)
    new_notebook = nbformat.v4.new_notebook()
    cell_list = []
    for slidenum, slide_contents in slide_dict.items():
        text_box = "\n".join(slide_contents["text_list"])
        verbatim = r"```"
        verbatim_box = [
            f"\n# Slide {slidenum}",
            f"{verbatim}\n",
            f"{text_box}\n",
            f"{verbatim}\n",
        ]
        if "image_list" in slide_contents:
            image_box = ["\n\n"]
            for image_file in slide_contents["image_list"]:
                image_box.append(f'\n<img src="{image_file}">\n')
            verbatim_box.append("\n".join(image_box))
        verbatim_string = "\n".join(verbatim_box)
        print(f"{verbatim_string}")
        new_cell = nbformat.v4.new_markdown_cell(source=verbatim_string)
        cell_list.append(new_cell)
    new_notebook["cells"] = cell_list
    markdown_file = pptx_file.with_suffix(".md")
    jupytext.write(new_notebook, markdown_file, fmt="md")


if __name__ == "__main__":
    main()
