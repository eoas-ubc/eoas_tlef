'''
conda install python-pptx
'''

from pathlib import Path
from pptx import Presentation
import nbformat
from jupytext.cli import jupytext

def get_text(slide):
    text_list = []
    for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_list.append(shape.text)
    return text_list


the_files = Path().glob("*.pptx")
for a_file in the_files:
    prs = Presentation(a_file)
    print(f"----- {a_file.name} ---------")
    for counter,slide in enumerate(prs.slides):
        print(f" ------ Slide {counter} ------")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
