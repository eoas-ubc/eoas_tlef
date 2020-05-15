# based on https://stackoverflow.com/questions/52491656/extracting-images-from-presentation-file

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
from PIL import Image
import click

class Extract_Pics:
    def __init__(self,filename,media_dir):
        """
        filename: name of pptx file

        """
        self.filename=Path(filename)
        self.image_num = 0
        self.deck_name = self.filename.stem
        self.slidenum = -1
        self.media_dir = Path(media_dir)
        
    def write_image(self,shape):
        """
        write the image in shape to a file, converting it
        to a png file if possible (this will work with jpg files,
        fail with wmf files)
        """
        image = shape.image
        # ---get image "file" contents---
        image_bytes = image.blob
        image_type = image.ext.lower()
        # ---make up a name for the file, e.g. 'image.jpg'---
        image_filename = f'{self.deck_name}-slide_{self.slidenum}-{self.image_num:03d}.{image_type}'
        self.image_num += 1
        full_path = self.media_dir / image_filename
        with open(full_path, 'wb') as f:
            print(f"writing {str(full_path)}")
            f.write(image_bytes)
        if image_type != 'png':
            png_path = full_path.with_suffix('.png')
            print(f"need to convert {image_filename}")
            try:
                with Image.open(full_path) as im:
                    im.load()
                    im.save(png_path)
                    print(f"saved converted file to {png_path}")
            except OSError as e:
                print(f"conversion failed {e}")

    def visitor(self,shape):
        """
        recursively find multiple shapes in a group and write
        them out to image files
        """
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"in a group")
            for count,s in enumerate(shape.shapes):
                print(f"group member {count}")
                self.visitor(s)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"individual picture")
            self.write_image(shape)

    def find_images(self):
        """
        find all group and individual images in a powerpoint
        presentations and write them out to media_dir
        """
        prs = Presentation(self.filename)
        for slidenum, slide in enumerate(prs.slides):
            self.slidenum=slidenum
            for shape in slide.shapes:
                self.visitor(shape)


@click.command()
@click.argument("pptx_file", type=str, nargs=1)
@click.argument("media_directory", type=str, nargs=1)
def main(pptx_file, media_directory):
    """\b
    usage: python image_extract.py pptx_file media_directory
    pptx_file: powerpoint file to convert
    media_directory: folder of png files to hold images from pptx
    """
    media_dir = Path(media_directory)
    media_dir.mkdir(parents=True,exist_ok=True)
    pptx_file = Path(pptx_file)
    do_job = Extract_Pics(pptx_file,media_dir)
    do_job.find_images()

                
if __name__ == "__main__":
    main()
